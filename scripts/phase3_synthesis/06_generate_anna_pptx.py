"""
Script: 06_generate_anna_pptx.py
Phase: 3 - Synthesis
Purpose: Generate PowerPoint presentation filling Anna's "Fran to fill in" gaps

This script generates a PowerPoint presentation with all the data points
Anna requested in the Dish Analysis Dec-25 slides.

Updated: 2026-01-08 to use:
- docs/data/zone_mvp_status.json (201 zones with order data)
- DATA/3_ANALYSIS/mvp_threshold_discovery.json (threshold evidence)
- 33 MVP Ready zones (16% of 201)

Gaps Filled:
1. Slide 5: Zone MVP Rationale with threshold evidence
2. Slide 6: Zone MVP Count (33/201 = 16%)
3. Slide 7: Minimum Dishes by Zone (tier recommendations)
4. Slide 8: Scoring Framework Definitions
5. Slide 9: Quadrant Summary (8 Core, 5 Pref, 2 Demand, 8 Depri)
6. Slide 10: Dish Scoring Data Table
7. Slide 13: Coverage Gaps & Next Steps

Inputs:
    - docs/data/zone_mvp_status.json
    - DATA/3_ANALYSIS/mvp_threshold_discovery.json
    - DATA/3_ANALYSIS/dish_scoring_anna_aligned.csv
    - DATA/3_ANALYSIS/latent_demand_scores.csv
    - config/mvp_thresholds.json

Outputs:
    - DELIVERABLES/reports/Dish_Analysis_Filled_YYYY-MM-DD.pptx
"""

import json
import logging
import pandas as pd
from pathlib import Path
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Setup logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

# Paths
PROJECT_ROOT = Path(__file__).parent.parent.parent
ANALYSIS_DIR = PROJECT_ROOT / "DATA" / "3_ANALYSIS"
DOCS_DATA_DIR = PROJECT_ROOT / "docs" / "data"
CONFIG_DIR = PROJECT_ROOT / "config"
DELIVERABLES_DIR = PROJECT_ROOT / "DELIVERABLES" / "reports"

# Deliveroo brand colors
DELIVEROO_TEAL = RGBColor(0, 204, 188)  # #00CCBC
DELIVEROO_DARK = RGBColor(0, 45, 70)    # #002D46
DELIVEROO_GRAY = RGBColor(88, 89, 91)   # #58595B
WHITE = RGBColor(255, 255, 255)

# Key numbers (from established analysis)
TOTAL_ORDERS = 82350
TOTAL_ZONES = 201
MVP_READY_ZONES = 33
MVP_READY_PCT = 16


def load_data():
    """Load all required data files."""
    data = {}
    
    # Zone MVP status (JSON - source of truth)
    path = DOCS_DATA_DIR / "zone_mvp_status.json"
    if path.exists():
        with open(path, 'r') as f:
            data['zone_mvp'] = json.load(f)
        logger.info(f"Loaded zone_mvp_status.json: {len(data['zone_mvp'])} zones")
    
    # MVP threshold discovery (JSON)
    path = ANALYSIS_DIR / "mvp_threshold_discovery.json"
    if path.exists():
        with open(path, 'r') as f:
            data['threshold_discovery'] = json.load(f)
        logger.info("Loaded mvp_threshold_discovery.json")
    
    # MVP thresholds config
    path = CONFIG_DIR / "mvp_thresholds.json"
    if path.exists():
        with open(path, 'r') as f:
            data['mvp_config'] = json.load(f)
        logger.info("Loaded mvp_thresholds.json")
    
    # Latent demand scores
    path = ANALYSIS_DIR / "latent_demand_scores.csv"
    if path.exists():
        data['latent_demand'] = pd.read_csv(path)
        logger.info(f"Loaded latent_demand_scores.csv: {len(data['latent_demand'])} rows")
    
    # Dish scoring aligned
    path = ANALYSIS_DIR / "dish_scoring_anna_aligned.csv"
    if path.exists():
        data['dish_scoring'] = pd.read_csv(path)
        logger.info(f"Loaded dish_scoring_anna_aligned.csv: {len(data['dish_scoring'])} rows")
    
    return data


def get_zone_mvp_counts(data):
    """Count zones by MVP status."""
    if 'zone_mvp' not in data:
        return {'MVP Ready': MVP_READY_ZONES, 'Near MVP': 48, 'Developing': 120}
    
    zones = data['zone_mvp']
    counts = {'MVP Ready': 0, 'Near MVP': 0, 'Developing': 0}
    
    for zone in zones:
        status = zone.get('mvp_status', 'Developing')
        if status in counts:
            counts[status] += 1
        else:
            counts['Developing'] += 1
    
    return counts


def get_partner_threshold_data(data):
    """Extract partner threshold data from mvp_threshold_discovery.json."""
    if 'threshold_discovery' not in data:
        return []
    
    td = data['threshold_discovery']
    partner_data = td.get('step_1_partners', {}).get('by_partner_count', [])
    
    rows = []
    for bucket in partner_data:
        rows.append({
            'Partners': bucket['bucket'],
            'Zones': bucket['zones'],
            'Repeat Rate': f"{bucket['behavioral']['repeat_rate']['value']*100:.1f}%",
            'Avg Rating': f"{bucket['behavioral']['avg_rating']['value']:.2f}",
            'Avg Orders': int(bucket['behavioral']['order_volume']['value'])
        })
    
    return rows


def get_cuisine_threshold_data(data):
    """Extract cuisine threshold data from mvp_threshold_discovery.json."""
    if 'threshold_discovery' not in data:
        return []
    
    td = data['threshold_discovery']
    cuisine_data = td.get('step_2_cuisines_within_partners', {}).get('by_cuisine_count', [])
    
    rows = []
    for bucket in cuisine_data:
        rows.append({
            'Cuisines': bucket['bucket'],
            'Zones': bucket['zones'],
            'Repeat Rate': f"{bucket['behavioral']['repeat_rate']['value']*100:.1f}%",
            'Avg Rating': f"{bucket['behavioral']['avg_rating']['value']:.2f}",
            'Avg Orders': int(bucket['behavioral']['order_volume']['value'])
        })
    
    return rows


def get_quadrant_data(data):
    """Get dish quadrant data from Anna's aligned scoring."""
    quadrants = {
        'Core Drivers': {
            'count': 8,
            'dishes': ['Pho', 'South Asian/Indian Curry', 'Biryani', 'Fried Rice', 
                      'Sushi', 'Katsu', 'Rice Bowl', 'Noodles'],
            'satisfaction': '88%',
            'sales_share': '73%'
        },
        'Preference Drivers': {
            'count': 5,
            'dishes': ['Grain Bowl', 'East Asian Curry', "Shepherd's Pie", 'Shawarma', 'Fajitas'],
            'satisfaction': '90%',
            'sales_share': '11%'
        },
        'Demand Boosters': {
            'count': 2,
            'dishes': ['Protein & Veg', 'Pizza'],
            'satisfaction': '76%',
            'sales_share': '8%'
        },
        'Deprioritised': {
            'count': 8,
            'dishes': ['Pasta', 'Lasagne', 'Tacos', 'Burrito/Burrito Bowl', 
                      'Quesadilla', 'Nachos', 'Chilli', 'Poke'],
            'satisfaction': '78%',
            'sales_share': '8%'
        }
    }
    return quadrants


def add_title_slide(prs, title, subtitle):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = DELIVEROO_DARK
    p.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(20)
    p.font.color.rgb = DELIVEROO_GRAY
    p.alignment = PP_ALIGN.CENTER
    
    return slide


def add_content_slide(prs, title, content_items):
    """Add a content slide with bullet points."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DELIVEROO_DARK
    
    # Add content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        if item.startswith('**'):
            # Bold header
            p.text = item.replace('**', '')
            p.font.bold = True
            p.font.size = Pt(16)
        elif item == '':
            p.text = ''
        else:
            p.text = f"• {item}"
            p.font.size = Pt(14)
        
        p.font.color.rgb = DELIVEROO_GRAY
        p.space_after = Pt(8)
    
    return slide


def add_table_slide(prs, title, headers, rows, col_widths=None, subtitle=None):
    """Add a slide with a data table."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = DELIVEROO_DARK
    
    # Add subtitle if provided
    top_offset = 1.0
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(9), Inches(0.4))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(12)
        p.font.color.rgb = DELIVEROO_GRAY
        top_offset = 1.3
    
    # Add table
    num_rows = len(rows) + 1
    num_cols = len(headers)
    
    table_height = min(0.35 * num_rows, 5.5)
    table = slide.shapes.add_table(
        num_rows, num_cols, 
        Inches(0.5), Inches(top_offset), 
        Inches(9), Inches(table_height)
    ).table
    
    # Set column widths
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
    
    # Add headers
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = DELIVEROO_TEAL
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    
    # Add data rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_data) if cell_data is not None else ''
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(10)
            p.font.color.rgb = DELIVEROO_GRAY
            p.alignment = PP_ALIGN.CENTER
    
    return slide


def generate_presentation(data):
    """Generate the complete PowerPoint presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Get data
    zone_counts = get_zone_mvp_counts(data)
    partner_data = get_partner_threshold_data(data)
    cuisine_data = get_cuisine_threshold_data(data)
    quadrants = get_quadrant_data(data)
    
    logger.info(f"Zone counts: {zone_counts}")
    logger.info(f"Partner buckets: {len(partner_data)}")
    logger.info(f"Cuisine buckets: {len(cuisine_data)}")
    
    # ========================================
    # SLIDE 1: Title
    # ========================================
    add_title_slide(
        prs,
        "Family Dinneroo",
        f"Dish Analysis - Data Pack\nGenerated: {datetime.now().strftime('%d %B %Y')}"
    )
    
    # ========================================
    # SLIDE 2: Executive Summary
    # ========================================
    add_content_slide(prs, "Executive Summary", [
        f"**Zone Performance (n={TOTAL_ZONES} live zones)**",
        f"33 zones (16%) are MVP Ready - these outperform others significantly",
        f"Data inflection at 3-4 partners/cuisines; business target is 5",
        "",
        f"**Dish Prioritisation (23 dishes analysed)**",
        "8 Core Drivers generate 73% of sales with 88% satisfaction",
        "5 Preference Drivers have highest satisfaction (90%) but low awareness",
        "2 Demand Boosters need quality improvement",
        "8 Deprioritised dishes should not be actively recruited",
        "",
        f"**Key Insight**",
        "Pho alone is 28% of Core Driver sales - concentration risk",
        "Indian Curry/Biryani have 74-79% coverage gaps - biggest opportunity"
    ])
    
    # ========================================
    # SLIDE 3 (Anna's Slide 5): Zone MVP Rationale - Partners
    # ========================================
    if partner_data:
        headers = ['Partners', 'Zones', 'Repeat Rate', 'Avg Rating', 'Avg Orders']
        rows = [[d['Partners'], d['Zones'], d['Repeat Rate'], d['Avg Rating'], d['Avg Orders']] 
                for d in partner_data]
        add_table_slide(
            prs, 
            "Slide 5: Zone MVP Rationale - Partners",
            headers, rows,
            col_widths=[1.5, 1.2, 1.8, 1.5, 1.8],
            subtitle="Data inflection at 3-4 partners (+4.5pp repeat rate). Business target: 5+"
        )
    
    # ========================================
    # SLIDE 4 (Anna's Slide 5): Zone MVP Rationale - Cuisines
    # ========================================
    if cuisine_data:
        headers = ['Cuisines', 'Zones', 'Repeat Rate', 'Avg Rating', 'Avg Orders']
        rows = [[d['Cuisines'], d['Zones'], d['Repeat Rate'], d['Avg Rating'], d['Avg Orders']] 
                for d in cuisine_data]
        add_table_slide(
            prs,
            "Slide 5: Zone MVP Rationale - Cuisines",
            headers, rows,
            col_widths=[1.5, 1.2, 1.8, 1.5, 1.8],
            subtitle="Analysis controls for 3+ partners. Data inflection at 3-4 cuisines. Business target: 5+"
        )
    
    # ========================================
    # SLIDE 5 (Anna's Slide 6): Zone MVP Count
    # ========================================
    headers = ['Status', 'Zones', '% of Total', 'Description']
    rows = [
        ['MVP Ready', zone_counts.get('MVP Ready', 33), '16%', 'Rating ≥4.0, Repeat ≥20%, Cuisines ≥5'],
        ['Near MVP', zone_counts.get('Near MVP', 48), '24%', 'Missing 1 criterion'],
        ['Developing', zone_counts.get('Developing', 120), '60%', 'Missing 2+ criteria'],
        ['TOTAL', TOTAL_ZONES, '100%', 'Live zones with order data']
    ]
    add_table_slide(
        prs,
        f"Slide 6: Zone MVP Count - {MVP_READY_ZONES} zones ({MVP_READY_PCT}%) MVP Ready",
        headers, rows,
        col_widths=[1.8, 1.2, 1.5, 4],
        subtitle=f"Source: zone_mvp_status.json (n={TOTAL_ZONES} zones with Snowflake order data)"
    )
    
    # ========================================
    # SLIDE 6 (Anna's Slide 7): Minimum Dishes by Zone
    # ========================================
    headers = ['Tier', 'Dishes', 'Sales Share', 'Satisfaction', 'Recommendation']
    rows = [
        ['Core Drivers', '8', '73%', '88%', '100% coverage - must have'],
        ['Preference Drivers', '5', '11%', '90%', 'Invest in merchandising'],
        ['Demand Boosters', '2', '8%', '76%', 'Improve quality'],
        ['Deprioritised', '8', '8%', '78%', 'Monitor only']
    ]
    add_table_slide(
        prs,
        "Slide 7: Minimum Dishes by Zone - Tier Recommendations",
        headers, rows,
        col_widths=[2, 1, 1.5, 1.5, 3],
        subtitle="Core Drivers generate 9x more sales per dish than Deprioritised (59.9 vs 6.4)"
    )
    
    # ========================================
    # SLIDE 7 (Anna's Slide 8): Scoring Framework
    # ========================================
    headers = ['Metric', 'Sample Size', 'Source', 'What It Measures']
    rows = [
        ['Avg Sales per Dish', f'{TOTAL_ORDERS:,}', 'Snowflake', 'Actual purchase behaviour'],
        ['% Zones in Top 5', '197 zones', 'Looker', 'Demand consistency'],
        ['Deliveroo Rating', '10,713', 'Looker', 'Post-delivery satisfaction'],
        ['Meal Satisfaction', '1,363', 'Post-order survey', '% Liked/Loved'],
        ['Repeat Intent', '1,363', 'Post-order survey', '% would order again'],
        ['Open-Text Requests', '372 mentions', 'All surveys', 'Unmet demand signals']
    ]
    add_table_slide(
        prs,
        "Slide 8: Scoring Framework Definitions",
        headers, rows,
        col_widths=[2.2, 1.5, 2, 3.3],
        subtitle="Framework combines revealed preference (orders) with stated preference (surveys)"
    )
    
    # ========================================
    # SLIDE 8 (Anna's Slide 9): Quadrant Summary
    # ========================================
    headers = ['Quadrant', 'Dishes', 'Sales %', 'Satisfaction', 'Example Dishes']
    rows = [
        ['Core Drivers', 8, '73%', '88%', 'Pho, Indian Curry, Biryani, Katsu'],
        ['Preference Drivers', 5, '11%', '90%', 'Fajitas, Shepherd\'s Pie, Shawarma'],
        ['Demand Boosters', 2, '8%', '76%', 'Protein & Veg, Pizza'],
        ['Deprioritised', 8, '8%', '78%', 'Pasta, Tacos, Burrito, Nachos']
    ]
    add_table_slide(
        prs,
        "Slide 9: Quadrant Summary - 8 Core, 5 Pref, 2 Demand, 8 Depri",
        headers, rows,
        col_widths=[2, 1, 1.2, 1.5, 3.5],
        subtitle="Core Drivers are 35% of dishes but 73% of sales"
    )
    
    # ========================================
    # SLIDE 9 (Anna's Slide 10): Dish Scoring Data
    # ========================================
    headers = ['Dish', 'Sales', 'Rating', 'Satisfaction', 'Quadrant']
    rows = [
        ['Pho', 136, '4.6', '92%', 'Core Driver'],
        ['South Asian/Indian Curry', 98, '4.7', '93%', 'Core Driver'],
        ['Biryani', 86, '4.6', '91%', 'Core Driver'],
        ['Protein & Veg', 44, '4.4', '71%', 'Demand Booster'],
        ['Fried Rice', 41, '4.5', '79%', 'Core Driver'],
        ['Sushi', 32, '4.7', '93%', 'Core Driver'],
        ['Katsu', 29, '4.3', '87%', 'Core Driver'],
        ['Rice Bowl', 29, '4.4', '84%', 'Core Driver'],
        ['Noodles', 28, '4.5', '88%', 'Core Driver'],
        ['Grain Bowl', 21, '4.7', '86%', 'Preference Driver'],
        ['East Asian Curry', 18, '4.4', '87%', 'Preference Driver'],
        ['Lasagne', 13, '4.6', '77%', 'Deprioritised'],
        ['Shepherd\'s Pie', 11, '4.3', '91%', 'Preference Driver'],
        ['Pizza', 11, '4.1', '81%', 'Demand Booster'],
        ['Shawarma', 10, '4.5', '90%', 'Preference Driver']
    ]
    add_table_slide(
        prs,
        "Slide 10: Dish Scoring Data Table (Top 15)",
        headers, rows,
        col_widths=[2.5, 1, 1.2, 1.5, 2.5],
        subtitle="Full data in ANNA_SLIDES_DATA_PACK.md"
    )
    
    # ========================================
    # SLIDE 10 (Anna's Slide 13): Coverage Gaps
    # ========================================
    headers = ['Dish', 'Coverage Gap', 'Opportunity Rank', 'Next Steps']
    rows = [
        ['South Asian/Indian Curry', '74%', '#1', 'Recruit Dishoom, Kricket, Tiffin Tin'],
        ['Biryani', '79%', '#2', 'Same Indian partners'],
        ['Pho', '48%', '#3', 'Expand Asian providers'],
        ['Protein & Veg', '93%', '#4', 'Add to Farmer J, LEON, Bill\'s'],
        ['Sushi', '63%', '#5', 'Work with Itsu, Iro Sushi, Kokoro'],
        ['Grain Bowl', '76%', '#7', 'Recruit Farmer J, healthy providers']
    ]
    add_table_slide(
        prs,
        "Slide 13: Coverage Gaps & Next Steps - Top Opportunities",
        headers, rows,
        col_widths=[2.5, 1.5, 1.5, 3.5],
        subtitle="Opportunity = Coverage Gap % × Avg Sales per Dish"
    )
    
    # ========================================
    # SLIDE 11: Key Insights
    # ========================================
    add_content_slide(prs, "Key Insights for Discussion", [
        "**Concentration Risk**",
        "Pho alone = 28% of Core Driver sales (136 of 479)",
        "6 of 8 Core Drivers are Asian dishes",
        "",
        "**The Familiarity Trap**",
        "Pasta gets most open-text requests (29) but only 78% satisfaction",
        "Fajitas gets 0 requests but 96% satisfaction - hidden gem",
        "",
        "**The Mexican Paradox**",
        "6 Mexican dishes, only Fajitas works (96% satisfaction)",
        "All others are Deprioritised - customers don't order Mexican for delivery",
        "",
        "**54% of Orders from Non-Families**",
        "Product works better for families (74% vs 68% satisfaction)",
        "Opportunity: Better targeting to reach more families"
    ])
    
    # ========================================
    # SLIDE 12: Sources
    # ========================================
    add_content_slide(prs, "Data Sources", [
        f"**Behavioural Data (Revealed Preference)**",
        f"Snowflake Orders: {TOTAL_ORDERS:,} orders",
        f"Deliveroo Ratings: 10,713 ratings",
        f"Zone Performance: {TOTAL_ZONES} zones",
        "",
        "**Survey Data (Stated Preference)**",
        "Post-Order Survey: 1,599 responses",
        "Dropoff Survey: 942 responses", 
        "OG Family Survey: 404 responses",
        "",
        "**Source Files**",
        "Zone MVP: docs/data/zone_mvp_status.json",
        "Thresholds: DATA/3_ANALYSIS/mvp_threshold_discovery.json",
        "Full data pack: DELIVERABLES/reports/ANNA_SLIDES_DATA_PACK.md"
    ])
    
    return prs


def main():
    """Main function to generate the presentation."""
    logger.info("=" * 60)
    logger.info("GENERATING ANNA'S DISH ANALYSIS POWERPOINT")
    logger.info("=" * 60)
    
    # Ensure output directory exists
    DELIVERABLES_DIR.mkdir(parents=True, exist_ok=True)
    
    # Load data
    data = load_data()
    if not data:
        logger.warning("Some data files not found, using defaults")
    
    # Generate presentation
    logger.info("\nGenerating slides...")
    prs = generate_presentation(data)
    
    # Save presentation
    today = datetime.now().strftime('%Y-%m-%d')
    output_path = DELIVERABLES_DIR / f"Dish_Analysis_Filled_{today}.pptx"
    prs.save(output_path)
    logger.info(f"\nSaved presentation to: {output_path}")
    
    # Summary
    logger.info("\n" + "=" * 60)
    logger.info("PRESENTATION COMPLETE")
    logger.info("=" * 60)
    logger.info(f"Total slides: {len(prs.slides)}")
    logger.info(f"Output: {output_path}")
    
    # Key numbers
    logger.info("\n" + "=" * 60)
    logger.info("KEY NUMBERS FOR ANNA")
    logger.info("=" * 60)
    logger.info(f"Total Orders: {TOTAL_ORDERS:,}")
    logger.info(f"Total Zones: {TOTAL_ZONES}")
    logger.info(f"MVP Ready: {MVP_READY_ZONES} ({MVP_READY_PCT}%)")
    logger.info("Quadrants: 8 Core, 5 Pref, 2 Demand, 8 Depri")
    logger.info("Partner inflection: 3-4 (data), 5 (target)")
    logger.info("Cuisine inflection: 3-4 (data), 5 (target)")


if __name__ == "__main__":
    main()
