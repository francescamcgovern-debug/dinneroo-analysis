#!/usr/bin/env python3
"""
Generate a "Cursor version" PPTX deck (template-free, clean build) for side-by-side comparison.

This deck is intentionally simple and text-forward: it mirrors the colleague deck's story flow
but uses Cursor canonical thresholds and scoring methodology references.

Inputs:
  - config/mvp_thresholds.json
  - config/scoring_framework_v3.json
  - config/dashboard_metrics.json (denominators)
  - docs/data/zone_mvp_status.json (established; do not recompute status labels)
  - DELIVERABLES/reports/deck_vs_cursor_diff_table.csv (for "what changed" appendix)

Output:
  - DELIVERABLES/presentations/cursor_version.pptx
"""

from __future__ import annotations

import csv
import json
from datetime import date
from pathlib import Path
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt


def _read_json(path: Path) -> Any:
    return json.loads(path.read_text(encoding="utf-8"))


def _read_csv(path: Path) -> List[Dict[str, str]]:
    with path.open("r", encoding="utf-8") as f:
        return list(csv.DictReader(f))


def _add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # title slide
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def _add_bullets_slide(prs: Presentation, title: str, bullets: List[str], notes: Optional[str] = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # title + content
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(20)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def _add_table_slide(prs: Presentation, title: str, headers: List[str], rows: List[List[str]], notes: Optional[str] = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # title only
    slide.shapes.title.text = title
    x, y, w, h = Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.2)
    table = slide.shapes.add_table(len(rows) + 1, len(headers), x, y, w, h).table

    for j, htxt in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = htxt
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(14)

    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)

    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def main() -> int:
    project_root = Path(__file__).resolve().parents[2]

    mvp_thresholds = _read_json(project_root / "config" / "mvp_thresholds.json")
    scoring = _read_json(project_root / "config" / "scoring_framework_v3.json")
    dashboard_metrics = _read_json(project_root / "config" / "dashboard_metrics.json")
    zone_mvp_status = _read_json(project_root / "docs" / "data" / "zone_mvp_status.json")
    diff_rows = _read_csv(project_root / "DELIVERABLES" / "reports" / "deck_vs_cursor_diff_table.csv")
    deck_zone_cov_path = project_root / "DELIVERABLES" / "presentation_data" / "deck_metrics_zone_threshold_coverage.csv"
    deck_zone_cov = _read_csv(deck_zone_cov_path) if deck_zone_cov_path.exists() else []

    out_dir = project_root / "DELIVERABLES" / "presentations"
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / "cursor_version.pptx"

    prs = Presentation()

    # Title
    _add_title_slide(
        prs,
        title="Family Dinneroo | Dish Analysis (Cursor version)",
        subtitle=f"Generated {date.today().isoformat()} from canonical configs + pipeline outputs",
    )

    # Summary (mirrors colleague structure)
    _add_bullets_slide(
        prs,
        "Summary",
        [
            "This deck is a side-by-side comparison version generated from Cursor source-of-truth.",
            "Zone MVP thresholds and scoring methodology are taken from config files (see notes).",
            "Where the colleague deck contains % coverage claims, this deck requires explicit denominators (total vs supply vs live zones).",
        ],
        notes=(
            "Sources:\n"
            "- config/mvp_thresholds.json\n"
            "- config/scoring_framework_v3.json\n"
            "- config/dashboard_metrics.json\n"
            "- docs/data/zone_mvp_status.json (established; do not recalc)\n"
        ),
    )

    # Gate A: MVP thresholds
    crit = mvp_thresholds.get("mvp_criteria", {})
    partners = crit.get("partners_min", {}).get("value")
    cuisines = crit.get("cuisines_min", {}).get("value")
    dishes = crit.get("dishes_min", {}).get("value")
    rating = crit.get("rating_min", {}).get("value")
    repeat = crit.get("repeat_rate_min", {}).get("value")

    _add_bullets_slide(
        prs,
        "Zone MVP (canonical)",
        [
            f"Partners: {partners}+",
            f"Cuisines (Core 7 count): {cuisines}+",
            f"Dishes: {dishes}+",
            f"Rating: {rating}+",
            f"Repeat rate: {repeat}%+",
        ],
        notes="Source: config/mvp_thresholds.json:mvp_criteria.* (business targets).",
    )

    # Denominators
    zm = dashboard_metrics.get("zone_metrics", {})
    total_zones = zm.get("total_zones_in_system", {}).get("value")
    supply_zones = zm.get("supply_zones_with_partners", {}).get("value")
    live_zones = zm.get("behavioral_zones_with_orders", {}).get("value")
    _add_bullets_slide(
        prs,
        "Zone denominators (must be explicit in % claims)",
        [
            f"Total zones (supply universe): {total_zones}",
            f"Supply zones with partners: {supply_zones}",
            f"Live zones with orders: {live_zones}",
        ],
        notes="Source: config/dashboard_metrics.json:zone_metrics.*",
    )

    if deck_zone_cov:
        # Present the live-zone denominator row as the most comparable to “active zones”
        live_row = next((r for r in deck_zone_cov if r.get("zone_denominator_type") == "live_zones_with_orders"), None)
        if live_row:
            _add_table_slide(
                prs,
                "Zone threshold coverage (explicit denominators)",
                headers=[
                    "Denominator",
                    "n zones",
                    "partners>=5",
                    "core7>=5",
                    "dishes>=21",
                    "all3",
                ],
                rows=[
                    [
                        live_row.get("zone_denominator_type", ""),
                        live_row.get("denominator", ""),
                        f"{live_row.get('pct_partners_min','')}%",
                        f"{live_row.get('pct_cuisines_min','')}%",
                        f"{live_row.get('pct_dishes_min','')}%",
                        f"{live_row.get('pct_all3','')}%",
                    ]
                ],
                notes=f"Source: {deck_zone_cov_path.relative_to(project_root)} (computed from established fields; not recalculating MVP status).",
            )

    # MVP status distribution from established file
    status_counts: Dict[str, int] = {}
    for z in zone_mvp_status:
        s = z.get("mvp_status", "UNKNOWN")
        status_counts[s] = status_counts.get(s, 0) + 1
    rows = [[k, str(v)] for k, v in sorted(status_counts.items(), key=lambda kv: kv[0])]
    _add_table_slide(
        prs,
        "Zone MVP status distribution (established)",
        headers=["Status", "Zones"],
        rows=rows,
        notes="Source: docs/data/zone_mvp_status.json (established status; do not recalc).",
    )

    # Gate B: Scoring methodology
    scoring_type = (scoring.get("scoring_method", {}) or {}).get("type")
    _add_bullets_slide(
        prs,
        "Dish scoring methodology (canonical)",
        [
            f"Framework: {scoring.get('description')}",
            f"Method: {scoring_type} (percentile-based 1–5 scoring)",
            "Outputs: separate lists for Family Performers, Couple Performers, Recruitment Priorities.",
        ],
        notes="Source: config/scoring_framework_v3.json",
    )

    # Factors per list
    lists = scoring.get("lists", {})
    factor_rows: List[List[str]] = []
    for list_name, spec in lists.items():
        factors = spec.get("factors", {})
        for factor, fmeta in factors.items():
            factor_rows.append([list_name, factor, str(fmeta.get("weight")), str(fmeta.get("source"))])
    _add_table_slide(
        prs,
        "Scoring factors by list (canonical)",
        headers=["List", "Factor", "Weight", "Source"],
        rows=factor_rows[:22],  # keep slide readable; full detail lives in config
        notes="Source: config/scoring_framework_v3.json:lists.*.factors",
    )

    # Appendix: key diffs (top DRIFT items)
    drift = [r for r in diff_rows if r.get("status") == "DRIFT"]
    drift_rows = []
    for r in drift[:12]:
        drift_rows.append([f"S{r.get('slide_index')}", r.get("deck_metric", ""), r.get("deck_text", "")[:40], r.get("cursor_source", "")])
    _add_table_slide(
        prs,
        "Appendix: Top deck ↔ Cursor drifts (auto-extracted)",
        headers=["Slide", "Metric", "Deck snippet", "Cursor source"],
        rows=drift_rows or [["-", "-", "-", "-"]],
        notes="Source: DELIVERABLES/reports/deck_vs_cursor_diff_table.csv",
    )

    prs.save(out_path)
    print(f"Wrote: {out_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

